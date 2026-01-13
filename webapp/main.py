"""
Slide Editor Web App - FastAPI Backend
交互式幻灯片编辑器后端
"""

import os
import sys
import uuid
import shutil
import tempfile
from pathlib import Path
from typing import Optional

# 加载 .env 文件
from dotenv import load_dotenv
load_dotenv()

from fastapi import FastAPI, File, UploadFile, HTTPException, Form
from fastapi.staticfiles import StaticFiles
from fastapi.responses import FileResponse, JSONResponse
from pydantic import BaseModel

# 添加父目录到路径以导入 slide_refiner
sys.path.insert(0, str(Path(__file__).parent.parent))

from PIL import Image
from pdf2image import convert_from_path
from google import genai
from google.genai import types
import img2pdf
import httpx
import re

app = FastAPI(title="Slide Editor", description="交互式幻灯片编辑器")

# 静态文件
app.mount("/static", StaticFiles(directory="static"), name="static")

# 会话存储
SESSIONS_DIR = Path("sessions")


def validate_session_id(session_id: str) -> bool:
    """验证 session_id 是否安全（防止路径遍历攻击）"""
    return bool(re.match(r'^[a-f0-9]{8}$', session_id))


def get_session_dir(session_id: str) -> Path:
    """安全获取会话目录"""
    if not validate_session_id(session_id):
        raise HTTPException(status_code=400, detail="无效的会话 ID")
    return SESSIONS_DIR / session_id
SESSIONS_DIR.mkdir(exist_ok=True)

# Gemini 客户端
def get_gemini_client():
    api_key = os.environ.get('GOOGLE_API_KEY') or os.environ.get('GEMINI_API_KEY')
    if not api_key:
        raise HTTPException(status_code=500, detail="请设置 GOOGLE_API_KEY 环境变量")
    
    custom_httpx_client = httpx.Client(http2=False, trust_env=False, timeout=600)
    http_options = types.HttpOptions(httpxClient=custom_httpx_client)
    return genai.Client(http_options=http_options)


class EditRequest(BaseModel):
    prompt: str


@app.get("/")
async def index():
    return FileResponse("static/index.html")


@app.post("/api/upload")
async def upload_pdf(file: UploadFile = File(...), remove_watermark: bool = Form(False), original_filename: Optional[str] = Form(None)):
    """上传 PDF 并转换为图片"""
    import traceback
    import logging
    
    logging.basicConfig(level=logging.INFO)
    logger = logging.getLogger(__name__)
    
    logger.info(f"收到上传请求: {file.filename}, 大小: {file.size if hasattr(file, 'size') else 'unknown'}")
    
    if not file.filename.endswith('.pdf'):
        raise HTTPException(status_code=400, detail="请上传 PDF 文件")
    
    try:
        # 创建会话
        session_id = str(uuid.uuid4())[:8]
        session_dir = SESSIONS_DIR / session_id
        session_dir.mkdir(exist_ok=True)
        logger.info(f"创建会话目录: {session_dir}")
        
        # 保存 PDF
        pdf_path = session_dir / "original.pdf"
        content = await file.read()
        logger.info(f"读取文件内容: {len(content)} 字节")
        
        with open(pdf_path, "wb") as f:
            f.write(content)
        logger.info(f"PDF 保存成功: {pdf_path}")
        
        # 转换为图片
        images_dir = session_dir / "original"
        images_dir.mkdir(exist_ok=True)
        
        try:
            images = convert_from_path(str(pdf_path), dpi=200)
            logger.info(f"PDF 转换成功: {len(images)} 页")
        except Exception as e:
            logger.error(f"PDF 转换失败: {e}")
            raise HTTPException(status_code=500, detail=f"PDF 转换失败: {str(e)}。请确保服务器已安装 poppler")
        
        pages = []
        for i, img in enumerate(images):
            img_path = images_dir / f"page_{i+1:03d}.png"
            img.save(str(img_path), "PNG")
            pages.append({
                "id": i + 1,
                "original": f"/api/sessions/{session_id}/original/{i+1}",
                "enhanced": None,
                "status": "pending"
            })
        
        # 创建增强图片目录
        (session_dir / "enhanced").mkdir(exist_ok=True)
        
        # 保存会话信息
        import json
        # 使用原始文件名（如果前端传了的话），否则用上传的文件名
        display_filename = original_filename or file.filename
        session_info = {
            "id": session_id,
            "filename": display_filename,
            "pages": pages,
            "remove_watermark": remove_watermark
        }
        with open(session_dir / "session.json", "w") as f:
            json.dump(session_info, f)
        
        logger.info(f"会话创建成功: {session_id}, 共 {len(pages)} 页")
        return {"session_id": session_id, "pages": pages, "total": len(pages)}
    
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"上传处理失败: {traceback.format_exc()}")
        raise HTTPException(status_code=500, detail=f"服务器处理失败: {str(e)}")


@app.get("/api/sessions/{session_id}")
async def get_session(session_id: str):
    """获取会话信息"""
    session_dir = get_session_dir(session_id)
    if not session_dir.exists():
        raise HTTPException(status_code=404, detail="会话不存在")
    
    import json
    with open(session_dir / "session.json") as f:
        return json.load(f)


@app.get("/api/sessions/{session_id}/original/{page_id}")
async def get_original_page(session_id: str, page_id: int):
    """获取原始页面图片"""
    session_dir = get_session_dir(session_id)
    img_path = session_dir / "original" / f"page_{page_id:03d}.png"
    if not img_path.exists():
        raise HTTPException(status_code=404, detail="页面不存在")
    return FileResponse(str(img_path), media_type="image/png")


@app.get("/api/sessions/{session_id}/enhanced/{page_id}")
async def get_enhanced_page(session_id: str, page_id: int):
    """获取增强后的页面图片"""
    session_dir = get_session_dir(session_id)
    img_path = session_dir / "enhanced" / f"page_{page_id:03d}.png"
    if not img_path.exists():
        raise HTTPException(status_code=404, detail="页面尚未增强")
    return FileResponse(str(img_path), media_type="image/png")


def detect_aspect_ratio(image: Image.Image) -> str:
    width, height = image.size
    ratio = width / height
    aspect_ratios = {
        "1:1": 1.0, "16:9": 16/9, "4:3": 4/3, "3:2": 3/2,
        "9:16": 9/16, "3:4": 3/4, "2:3": 2/3
    }
    closest = min(aspect_ratios.items(), key=lambda x: abs(x[1] - ratio))
    return closest[0]


def blank_watermark_area(image: Image.Image) -> Image.Image:
    """移除右下角水印区域"""
    width, height = image.size
    corner_width = int(width * 0.15)
    corner_height = int(height * 0.08)
    x1, y1 = width - corner_width, height - corner_height
    
    sample_y = max(0, y1 - 10)
    sample_x = x1 + corner_width // 2
    try:
        bg_color = image.getpixel((sample_x, sample_y))
    except:
        bg_color = (245, 245, 245)
    
    result = image.copy()
    from PIL import ImageDraw
    draw = ImageDraw.Draw(result)
    draw.rectangle([x1, y1, width, height], fill=bg_color)
    return result


@app.post("/api/sessions/{session_id}/enhance/{page_id}")
async def enhance_page(session_id: str, page_id: int, custom_prompt: Optional[str] = Form(None), remove_watermark: bool = Form(False)):
    """增强单个页面"""
    session_dir = get_session_dir(session_id)
    if not session_dir.exists():
        raise HTTPException(status_code=404, detail="会话不存在")
    
    import json
    with open(session_dir / "session.json") as f:
        session = json.load(f)
    
    original_path = session_dir / "original" / f"page_{page_id:03d}.png"
    enhanced_path = session_dir / "enhanced" / f"page_{page_id:03d}.png"
    
    if not original_path.exists():
        raise HTTPException(status_code=404, detail="原始页面不存在")
    
    # 构建 prompt
    if custom_prompt:
        prompt = f"""Enhance this presentation slide based on the following instructions:

{custom_prompt}

Also apply these enhancements:
1. SHARPEN all text to be crisp and highly readable
2. ENHANCE image quality - reduce blur and noise
3. IMPROVE color vibrancy
4. OUTPUT at maximum resolution"""
    else:
        base_prompt = """Enhance this presentation slide to ultra-high definition quality.

CRITICAL RULES:
1. PRESERVE all content exactly
2. SHARPEN all text to be crisp and highly readable
3. ENHANCE image quality - reduce blur, noise, and compression artifacts
4. IMPROVE color vibrancy while maintaining the original color scheme
5. OUTPUT at maximum resolution"""
        
        if remove_watermark:
            prompt = base_prompt + """
6. IMPORTANT: There is a BLANK area in the bottom-right corner. Fill it seamlessly with the surrounding background.

This is an image enhancement and inpainting task."""
        else:
            prompt = base_prompt + "\n\nThis is ONLY an image quality enhancement task."
    
    # 加载图片
    image = Image.open(str(original_path))
    
    # 如果需要移除水印
    if remove_watermark:
        image = blank_watermark_area(image)
    
    aspect_ratio = detect_aspect_ratio(image)
    
    # 调用 Gemini
    try:
        client = get_gemini_client()
        response = client.models.generate_content(
            model="gemini-3-pro-image-preview",
            contents=[prompt, image],
            config=types.GenerateContentConfig(
                response_modalities=['TEXT', 'IMAGE'],
                image_config=types.ImageConfig(
                    aspect_ratio=aspect_ratio,
                    image_size="4K"
                ),
            )
        )
        
        if response is None or response.parts is None:
            raise HTTPException(status_code=500, detail="API 返回空响应")
        
        saved = False
        for part in response.parts:
            if hasattr(part, 'thought') and part.thought:
                continue
            if part.inline_data is not None:
                enhanced_image = part.as_image()
                enhanced_image.save(str(enhanced_path))
                saved = True
        
        if not saved:
            raise HTTPException(status_code=500, detail="未返回增强图片")
        
        # 更新会话
        for page in session["pages"]:
            if page["id"] == page_id:
                page["enhanced"] = f"/api/sessions/{session_id}/enhanced/{page_id}"
                page["status"] = "done"
                break
        
        with open(session_dir / "session.json", "w") as f:
            json.dump(session, f)
        
        return {"success": True, "enhanced": f"/api/sessions/{session_id}/enhanced/{page_id}"}
    
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@app.post("/api/sessions/{session_id}/enhance-all")
async def enhance_all_pages(session_id: str):
    """增强所有页面"""
    session_dir = get_session_dir(session_id)
    if not session_dir.exists():
        raise HTTPException(status_code=404, detail="会话不存在")
    
    import json
    with open(session_dir / "session.json") as f:
        session = json.load(f)
    
    results = []
    for page in session["pages"]:
        if page["status"] != "done":
            try:
                result = await enhance_page(session_id, page["id"])
                results.append({"id": page["id"], "success": True})
            except Exception as e:
                results.append({"id": page["id"], "success": False, "error": str(e)})
    
    return {"results": results}


@app.post("/api/sessions/{session_id}/export/pdf")
async def export_pdf(session_id: str):
    """导出为 PDF"""
    session_dir = get_session_dir(session_id)
    if not session_dir.exists():
        raise HTTPException(status_code=404, detail="会话不存在")
    
    import json
    with open(session_dir / "session.json") as f:
        session = json.load(f)
    
    # 收集所有图片（优先使用增强版）
    image_paths = []
    for page in session["pages"]:
        enhanced = session_dir / "enhanced" / f"page_{page['id']:03d}.png"
        original = session_dir / "original" / f"page_{page['id']:03d}.png"
        image_paths.append(str(enhanced if enhanced.exists() else original))
    
    output_path = session_dir / "output.pdf"
    with open(output_path, "wb") as f:
        f.write(img2pdf.convert(image_paths))
    
    # 安全提取原始文件名
    original_name = Path(session['filename']).stem
    return FileResponse(str(output_path), filename=f"{original_name}_enhanced.pdf")


@app.post("/api/sessions/{session_id}/export/pptx")
async def export_pptx(session_id: str):
    """导出为 PPTX"""
    from pptx import Presentation
    from pptx.util import Inches
    
    session_dir = get_session_dir(session_id)
    if not session_dir.exists():
        raise HTTPException(status_code=404, detail="会话不存在")
    
    import json
    with open(session_dir / "session.json") as f:
        session = json.load(f)
    
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    for page in session["pages"]:
        enhanced = session_dir / "enhanced" / f"page_{page['id']:03d}.png"
        original = session_dir / "original" / f"page_{page['id']:03d}.png"
        img_path = str(enhanced if enhanced.exists() else original)
        
        blank_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_layout)
        slide.shapes.add_picture(img_path, Inches(0), Inches(0), 
                                  width=prs.slide_width, height=prs.slide_height)
    
    output_path = session_dir / "output.pptx"
    prs.save(str(output_path))
    
    # 安全提取原始文件名
    original_name = Path(session['filename']).stem
    return FileResponse(str(output_path), filename=f"{original_name}_enhanced.pptx")


@app.post("/api/sessions/{session_id}/reset/{page_id}")
async def reset_page(session_id: str, page_id: int):
    """重置单个页面（删除增强版本）"""
    session_dir = get_session_dir(session_id)
    if not session_dir.exists():
        raise HTTPException(status_code=404, detail="会话不存在")
    
    import json
    with open(session_dir / "session.json") as f:
        session = json.load(f)
    
    # 删除增强后的图片
    enhanced_path = session_dir / "enhanced" / f"page_{page_id:03d}.png"
    if enhanced_path.exists():
        enhanced_path.unlink()
    
    # 更新会话状态
    for page in session["pages"]:
        if page["id"] == page_id:
            page["enhanced"] = None
            page["status"] = "pending"
            break
    
    with open(session_dir / "session.json", "w") as f:
        json.dump(session, f)
    
    return {"success": True, "original": f"/api/sessions/{session_id}/original/{page_id}"}


@app.delete("/api/sessions/{session_id}")
async def delete_session(session_id: str):
    """删除会话"""
    session_dir = get_session_dir(session_id)
    if session_dir.exists():
        shutil.rmtree(session_dir)
    return {"success": True}


if __name__ == "__main__":
    import uvicorn
    uvicorn.run(app, host="0.0.0.0", port=9998)
